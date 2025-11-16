"""
AI-powered document analysis service
Extracts text from documents and analyzes content using Google Gemini
"""
import os
import PyPDF2
import docx
from pptx import Presentation
import google.generativeai as genai
from typing import Dict, Optional
import re
from app.config import settings

class DocumentAnalyzer:
    def __init__(self):
        api_key = settings.GEMINI_API_KEY
        if api_key:
            genai.configure(api_key=api_key)
            self.model = genai.GenerativeModel('gemini-2.5-flash')
        else:
            self.model = None
            print("Warning: GEMINI_API_KEY not set. AI analysis will be disabled.")
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return ""
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text based on file type"""
        file_type = file_type.lower()
        
        if file_type == 'pdf' or file_path.endswith('.pdf'):
            return self.extract_text_from_pdf(file_path)
        elif file_type in ['docx', 'doc'] or file_path.endswith(('.docx', '.doc')):
            return self.extract_text_from_docx(file_path)
        elif file_type in ['pptx', 'ppt'] or file_path.endswith(('.pptx', '.ppt')):
            return self.extract_text_from_pptx(file_path)
        else:
            return ""
    
    def analyze_document_content(self, text: str, filename: str) -> Dict:
        """Analyze document content using AI"""
        if not self.model or not text:
            # Fallback to basic filename analysis
            return self._basic_analysis(filename, text)
        
        try:
            # Limit text to first 3000 characters for analysis
            text_sample = text[:3000] if len(text) > 3000 else text
            
            prompt = f"""Analyze this educational document and provide a JSON response with the following fields:

Document Filename: {filename}
Document Content:
{text_sample}

Provide analysis in this exact JSON format:
{{
    "document_type": "one of: syllabus, assignment, notes, question_paper, tutorial, project, presentation, notice, general",
    "subject": "the academic subject (e.g., Mathematics, Computer Science, Physics)",
    "semester": "semester number if mentioned (1-8), or null",
    "branch": "branch/department if mentioned (e.g., CSE, ISE, ECE, Mechanical), or null",
    "topics": ["list of 3-5 main topics covered"],
    "keywords": ["list of 5-8 relevant keywords"],
    "description": "2-3 sentence summary of the document"
}}

Be precise and return only valid JSON."""

            response = self.model.generate_content(prompt)
            result_text = response.text.strip()
            
            # Extract JSON from response (handle markdown code blocks)
            json_match = re.search(r'\{[\s\S]*\}', result_text)
            if json_match:
                import json
                analysis = json.loads(json_match.group())
                
                # Validate and normalize
                return {
                    'document_type': analysis.get('document_type', 'general'),
                    'subject': analysis.get('subject'),
                    'semester': analysis.get('semester'),
                    'branch': analysis.get('branch'),
                    'topics': analysis.get('topics', []),
                    'keywords': analysis.get('keywords', []),
                    'description': analysis.get('description', ''),
                    'ai_analyzed': True
                }
            else:
                return self._basic_analysis(filename, text)
                
        except Exception as e:
            print(f"AI analysis error: {e}")
            return self._basic_analysis(filename, text)
    
    def _basic_analysis(self, filename: str, text: str) -> Dict:
        """Fallback basic analysis based on filename and keywords"""
        filename_lower = filename.lower()
        text_lower = text.lower() if text else ""
        
        # Detect document type
        doc_type = 'general'
        if any(word in filename_lower for word in ['syllabus', 'curriculum']):
            doc_type = 'syllabus'
        elif any(word in filename_lower for word in ['assignment', 'homework', 'task']):
            doc_type = 'assignment'
        elif any(word in filename_lower for word in ['notes', 'lecture', 'chapter']):
            doc_type = 'notes'
        elif any(word in filename_lower for word in ['question', 'exam', 'test', 'quiz']):
            doc_type = 'question_paper'
        elif any(word in filename_lower for word in ['tutorial', 'lab', 'practical']):
            doc_type = 'tutorial'
        elif any(word in filename_lower for word in ['project', 'proposal']):
            doc_type = 'project'
        elif any(word in filename_lower for word in ['presentation', 'slides', 'ppt']):
            doc_type = 'presentation'
        elif any(word in filename_lower for word in ['notice', 'circular', 'announcement']):
            doc_type = 'notice'
        
        # Detect semester
        semester = None
        sem_match = re.search(r'sem(?:ester)?[\s_-]?(\d)', filename_lower)
        if sem_match:
            semester = int(sem_match.group(1))
        
        # Detect subject (basic keywords)
        subject = None
        subjects = {
            'mathematics': ['math', 'mathematics', 'calculus', 'algebra'],
            'physics': ['physics', 'mechanics', 'thermodynamics'],
            'chemistry': ['chemistry', 'organic', 'inorganic'],
            'computer science': ['computer', 'programming', 'algorithm', 'data structure', 'python', 'java'],
            'electronics': ['electronics', 'circuits', 'vlsi', 'embedded'],
            'mechanical': ['mechanical', 'thermodynamics', 'fluid'],
        }
        
        for subj, keywords in subjects.items():
            if any(kw in filename_lower or kw in text_lower for kw in keywords):
                subject = subj.title()
                break
        
        return {
            'document_type': doc_type,
            'subject': subject,
            'semester': semester,
            'branch': None,
            'topics': [],
            'keywords': [],
            'description': f"Document: {filename}",
            'ai_analyzed': False
        }
    
    def analyze_document(self, file_path: str, filename: str, file_type: str) -> Dict:
        """Complete document analysis pipeline"""
        # Extract text
        text = self.extract_text(file_path, file_type)
        
        # Analyze content
        analysis = self.analyze_document_content(text, filename)
        
        # Add extracted text (truncated for storage)
        analysis['extracted_text'] = text[:1000] if text else ""
        analysis['text_length'] = len(text)
        
        return analysis
